from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os

pptx_path = r"D:\Opencode files\html簡報技能\“嘉茵頤縫”一次性可吸收釘皮內吻合器-敏盛開刀房.pptx"

prs = Presentation(pptx_path)

data = {
    'slides': [],
    'slide_width': prs.slide_width,
    'slide_height': prs.slide_height
}

for i, slide in enumerate(prs.slides):
    slide_data = {
        'slide_number': i + 1,
        'shapes': []
    }
    for shape in slide.shapes:
        shape_data = {
            'type': shape.shape_type,
            'name': shape.name,
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height
        }
        if shape.has_text_frame:
            text = ''
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text += run.text
                text += '\n'
            shape_data['text'] = text.strip()
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data.append(row_data)
            shape_data['table'] = table_data
        slide_data['shapes'].append(shape_data)
    data['slides'].append(slide_data)

with open('slides_data.json', 'w', encoding='utf-8') as f:
    json.dump(data, f, ensure_ascii=False, indent=2)

print(f'Total slides: {len(prs.slides)}')
print('Data saved to slides_data.json')